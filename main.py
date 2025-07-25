from fastapi import FastAPI, Form
from fastapi.responses import FileResponse, HTMLResponse
from pptx import Presentation
import os, zipfile

app = FastAPI()

@app.get("/", response_class=HTMLResponse)
def form():
    return """
    <h2>Введите имя или несколько имён через запятую</h2>
    <form action="/generate" method="post">
        <input type="text" name="name" placeholder="Арман, Ержан" style="padding: 8px; width: 250px;">
        <button type="submit" style="padding: 8px 12px; margin-left: 10px;">Создать кувертки</button>
    </form>
    """

@app.post("/generate")
def generate(name: str = Form(...)):
    names = [n.strip() for n in name.split(",") if n.strip()]
    os.makedirs("output", exist_ok=True)
    output_files = []

    for n in names:
        prs = Presentation("template.pptx")
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "ИМЯ" in run.text:
                                run.text = run.text.replace("ИМЯ", n)
        filename = f"output/{n}_кувертка.pptx"
        prs.save(filename)
        output_files.append(filename)

    if len(output_files) == 1:
        return FileResponse(output_files[0], filename=os.path.basename(output_files[0]))
    else:
        zip_path = "output/кувертки.zip"
        with zipfile.ZipFile(zip_path, "w") as zipf:
            for f in output_files:
                zipf.write(f, arcname=os.path.basename(f))
        return FileResponse(zip_path, filename="кувертки.zip")
